"""
storyteller/pptx_service.py

Service layer for generating rich, dynamically styled PowerPoint (.pptx) files
from Fafanua presentation models.
"""

from io import BytesIO
from typing import Any, Dict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class PPTXGenerator:
    """Generates PowerPoint (.pptx) presentations with theme-based design styling."""

    THEMES: Dict[str, Dict[str, RGBColor]] = {
        "bold": {
            "background": RGBColor(15, 23, 42),       # Slate 900
            "badge_bg": RGBColor(30, 41, 59),         # Slate 800
            "badge_text": RGBColor(147, 197, 253),     # Blue 300
            "title_text": RGBColor(248, 250, 252),     # Slate 50
            "bullet_text": RGBColor(203, 213, 225),    # Slate 300
            "bullet_dot": RGBColor(96, 165, 250),      # Blue 400
            "container_bg": RGBColor(30, 41, 59),     # Slate 800
            "container_border": RGBColor(51, 65, 85), # Slate 700
            "code_text": RGBColor(148, 163, 184),      # Slate 400
        },
        "warm": {
            "background": RGBColor(44, 24, 16),        # Amber 950
            "badge_bg": RGBColor(69, 26, 3),          # Amber 900
            "badge_text": RGBColor(252, 211, 77),      # Amber 300
            "title_text": RGBColor(254, 243, 199),      # Amber 100
            "bullet_text": RGBColor(245, 208, 140),     # Amber 200
            "bullet_dot": RGBColor(251, 191, 36),       # Amber 400
            "container_bg": RGBColor(69, 26, 3),      # Amber 900
            "container_border": RGBColor(120, 53, 15),  # Amber 800
            "code_text": RGBColor(217, 119, 6),        # Amber 600
        },
        "clean": {
            "background": RGBColor(6, 78, 59),         # Emerald 950
            "badge_bg": RGBColor(6, 95, 70),          # Emerald 900
            "badge_text": RGBColor(110, 231, 183),     # Emerald 300
            "title_text": RGBColor(209, 250, 229),     # Emerald 100
            "bullet_text": RGBColor(167, 243, 208),    # Emerald 200
            "bullet_dot": RGBColor(52, 211, 153),      # Emerald 400
            "container_bg": RGBColor(6, 95, 70),      # Emerald 900
            "container_border": RGBColor(4, 120, 87),   # Emerald 800
            "code_text": RGBColor(52, 211, 153),       # Emerald 400
        },
        "analytical": {
            "background": RGBColor(58, 12, 92),        # Purple 950
            "badge_bg": RGBColor(76, 29, 149),        # Purple 900
            "badge_text": RGBColor(216, 180, 254),     # Purple 300
            "title_text": RGBColor(243, 232, 255),     # Purple 100
            "bullet_text": RGBColor(233, 213, 255),    # Purple 200
            "bullet_dot": RGBColor(192, 132, 252),      # Purple 400
            "container_bg": RGBColor(76, 29, 149),     # Purple 900
            "container_border": RGBColor(109, 40, 217),# Purple 700
            "code_text": RGBColor(192, 132, 252),      # Purple 400
        },
    }

    def generate_presentation(self, presentation_obj: Any) -> BytesIO:
        """Generates a PowerPoint (.pptx) byte stream for a presentation object.

        Args:
            presentation_obj: Django Presentation instance with related slides.

        Returns:
            BytesIO stream containing binary PPTX data.
        """
        prs = Presentation()
        # Set 16:9 Widescreen aspect ratio (13.333 x 7.5 inches)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # Blank layout

        slides = list(presentation_obj.slides.all())

        for idx, slide_model in enumerate(slides):
            slide = prs.slides.add_slide(blank_layout)

            # Resolve theme colors with 'clean' as fallback
            variant = (slide_model.theme_variant or "clean").lower()
            theme = self.THEMES.get(variant, self.THEMES["clean"])

            # 1. Slide Background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = theme["background"]

            # 2. Header / Badge Row
            header_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.4)
            )
            tf_header = header_box.text_frame
            tf_header.word_wrap = True
            p_hdr = tf_header.paragraphs[0]
            p_hdr.text = f"SLIDE {idx + 1}  •  {variant.upper()}"
            p_hdr.font.size = Pt(11)
            p_hdr.font.bold = True
            p_hdr.font.name = "Arial"
            p_hdr.font.color.rgb = theme["badge_text"]

            # 3. Slide Title
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.9), Inches(11.733), Inches(0.9)
            )
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = slide_model.title
            p_title.font.size = Pt(26)
            p_title.font.bold = True
            p_title.font.name = "Arial"
            p_title.font.color.rgb = theme["title_text"]

            # 4. Bullet Points
            has_diagram = bool(
                slide_model.diagram_code and slide_model.diagram_code.strip()
            )
            bullet_height = Inches(3.2) if has_diagram else Inches(5.0)

            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.9), Inches(11.733), bullet_height
            )
            tf_content = content_box.text_frame
            tf_content.word_wrap = True

            bullets = slide_model.bullet_points or []
            for b_idx, bullet_text in enumerate(bullets):
                p = (
                    tf_content.paragraphs[0]
                    if b_idx == 0
                    else tf_content.add_paragraph()
                )
                p.space_after = Pt(10)

                # Bullet symbol run
                run_dot = p.add_run()
                run_dot.text = "•  "
                run_dot.font.size = Pt(14)
                run_dot.font.name = "Arial"
                run_dot.font.bold = True
                run_dot.font.color.rgb = theme["bullet_dot"]

                # Bullet text run
                run_text = p.add_run()
                run_text.text = str(bullet_text)
                run_text.font.size = Pt(14)
                run_text.font.name = "Arial"
                run_text.font.color.rgb = theme["bullet_text"]

            # 5. Diagram Fallback Box (if slide contains Mermaid syntax)
            if has_diagram:
                box_top = Inches(5.2)
                box_height = Inches(1.7)
                box_width = Inches(11.733)

                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.8),
                    box_top,
                    box_width,
                    box_height,
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = theme["container_bg"]
                shape.line.color.rgb = theme["container_border"]
                shape.line.width = Pt(1)

                tf_diag = shape.text_frame
                tf_diag.word_wrap = True
                tf_diag.margin_left = Inches(0.2)
                tf_diag.margin_top = Inches(0.15)
                tf_diag.margin_right = Inches(0.2)
                tf_diag.margin_bottom = Inches(0.15)

                p_diag_lbl = tf_diag.paragraphs[0]
                p_diag_lbl.text = (
                    "[Contains Mermaid Diagram: View in Web App or paste code into mermaid.live]"
                )
                p_diag_lbl.font.size = Pt(10)
                p_diag_lbl.font.bold = True
                p_diag_lbl.font.name = "Arial"
                p_diag_lbl.font.color.rgb = theme["badge_text"]
                p_diag_lbl.space_after = Pt(4)

                p_diag_code = tf_diag.add_paragraph()
                p_diag_code.text = slide_model.diagram_code.strip()
                p_diag_code.font.size = Pt(9)
                p_diag_code.font.name = "Courier New"
                p_diag_code.font.color.rgb = theme["code_text"]

        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
